"""
Update AccessBridge_Presentation.pptx for Session 8 submission polish.

Surgical edits only — replaces stale numbers/text in place, appends 2 new slides
(Roadmap, QA Summary). Does not rebuild the deck from scratch.

Inputs: AccessBridge_Presentation.pptx  (read-only source)
Output: AccessBridge_Presentation_v2.pptx
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor as _RGBColor


# ---------------------------------------------------------------------------
# Text replacements per slide (1-indexed). Matches exact text; order matters.
# ---------------------------------------------------------------------------

REPLACEMENTS = {
    1: [
        ("10+", "28"),
        ("25+", "45+"),
        ("116", "544"),
    ],
    4: [
        (
            "Build: Vite + TypeScript strict, output 132KB content + 28KB background + 19KB sidepanel",
            "Build: Vite + TypeScript strict. Output: 322 KB content, 36 KB background, 30 KB popup, 414 KB sidepanel",
        ),
    ],
    7: [
        ("Voice Commands (25+)", "Voice Commands (45+ EN/Hindi, 22 Indic languages)"),
    ],
    9: [
        (
            "Extensible Architecture",
            "Six Shipping Connectors",
        ),
        (
            "DomainConnectorRegistry auto-detects domain from URL → activates matching connector → injects domain-specific UI + helpers. Add new connectors in <50 lines.",
            "Banking · Insurance · Healthcare (drug interactions) · Telecom (bill-shock alerts) · Retail (savings badges) · Manufacturing (hazard keywords). DomainConnectorRegistry routes by hostname; add new connectors in <50 lines.",
        ),
    ],
    11: [
        ("116 tests, zero errors, strict TypeScript, automated deployment pipeline",
         "544 tests, zero errors, strict TypeScript, automated deployment pipeline"),
        ("116", "544"),
        ("3", "14"),
        ("Test Suites", "Test Files"),
        (
            "StruggleDetector: 16 tests (signal processing, baseline, scoring) | DecisionEngine: 21 tests (rules, confidence, revert, profiles) | ProfileStore: 25 tests (CRUD, encryption, migration, edge cases) | AI Engine: 54 tests (cache, normal",
            "Core (382 tests): StruggleDetector 18 · DecisionEngine 24 · ProfileStore 20 · Audit engine + rules 96 · Profile versioning + drift 24 · Environment detect 38 · Gesture recognizer + bindings 36 · Indic language-ranges 30 · Language detect 28 · Shortcut DSL 19 · Transliteration 49. AI Engine (54 tests): cache, normalizer, cost tracker, local provider. Extension (108 tests): captions, observatory publisher, deepenings, indic-commands, env sensor, gestures, time-awareness, action-items.",
        ),
    ],
    13: [
        ("Team AccessBridge", "Manish Kumar"),
        ("25+ English + Hindi", "45+ EN + Hindi + 20 Indic"),
        ("Banking + Insurance", "6 verticals (Banking, Insurance, Healthcare, Telecom, Retail, Manufacturing)"),
    ],
}


# ---------------------------------------------------------------------------
# Roadmap + QA Summary slide content. Both use same dark surface layout.
# ---------------------------------------------------------------------------

ROADMAP = {
    "title": "Roadmap",
    "subtitle": "Phase 1 shipped — Phase 2 + 3 unlock B2B scale",
    "phases": [
        (
            "Phase 1  ·  Shipped in this submission",
            "Chrome Manifest V3 extension feature-complete.\n3 modules × 28 features.\n6 domain connectors with v1 deepening.\n22 Indian languages (full Web Speech + DOM vocab).\nCompliance observatory live with differential privacy + Merkle commits.\nSingle monorepo (@accessbridge/core + ai-engine + extension) ready for reuse.",
        ),
        (
            "Phase 2  ·  Next 12 weeks",
            "Desktop companion (Tauri) — overlay on native Word/Teams/VS Code.\nProfile sync + SSO — settings follow user across devices.\nOn-device ONNX models — Whisper STT + Transformers.js so local tier handles the 20% now escalating to cloud.\nAndroid AccessibilityService prototype for native apps.",
        ),
        (
            "Phase 3  ·  6 months",
            "iOS Safari extension + in-app SDK.\nEnterprise admin console (SCCM/Intune silent install, ROI reporting per app).\nPublic JS/TS SDK + REST API for third-party web app embedding.\nCommunity-contributed domain connectors.",
        ),
    ],
}

QA_SUMMARY_PLACEHOLDER = {
    "title": "Chrome Sideload QA Summary",
    "subtitle": "Real browser, real sites, honest pass/fail — full report in QA_REPORT.md",
    # These get filled by user at session end; leave placeholders callable by a
    # later pass once QA counts are known.
    "stats": [
        ("54", "Tests"),
        ("__", "Pass"),
        ("__", "Partial"),
        ("__", "Fail"),
    ],
    "notes": (
        "Popup UI (28) · Content on live sites (20) · Side panel (8) · "
        "Background SW (8) · VPS integration (3) · Error recovery (5).\n\n"
        "Every FAIL has a BUG-XXX entry in RCA.md with the fix commit. "
        "Anything shipping with known issues is disclosed in QA_REPORT.md § Known Issues."
    ),
}


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------


def iter_text_frames(slide):
    for sh in slide.shapes:
        if sh.has_text_frame:
            yield sh.text_frame


def replace_in_slide(slide, pairs):
    """Replace each (old, new) pair in the first matching text frame.

    Matching rule: compare the aggregated text of a frame. If old == aggregated
    text exactly, replace that frame's first run text with new (clearing others).
    Fall back to per-run substring replacement otherwise.
    """
    for old, new in pairs:
        changed = False
        for tf in iter_text_frames(slide):
            agg = tf.text
            # Exact-match replacement (most surgical)
            if agg.strip() == old.strip():
                for para in tf.paragraphs:
                    for run in para.runs:
                        run.text = ""
                if tf.paragraphs and tf.paragraphs[0].runs:
                    tf.paragraphs[0].runs[0].text = new
                else:
                    tf.text = new
                changed = True
                break
            # Substring replacement within a single run (preserves formatting)
            for para in tf.paragraphs:
                for run in para.runs:
                    if old in run.text:
                        run.text = run.text.replace(old, new)
                        changed = True
                        break
                if changed:
                    break
            if changed:
                break
        if not changed:
            print(f"  [warn] not found: {old!r}")


def add_dark_slide(prs, title, subtitle):
    """Add a blank-layout slide with dark-theme title + subtitle textboxes."""
    layout = prs.slide_layouts[6]  # blank
    s = prs.slides.add_slide(layout)

    # Dark background (matches existing deck)
    from pptx.oxml.ns import qn
    from lxml import etree
    bg = s.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = _RGBColor(0x0a, 0x0a, 0x1a)

    title_box = s.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12.33), Inches(0.9))
    p = title_box.text_frame.paragraphs[0]
    p.text = title
    p.runs[0].font.size = Pt(36)
    p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = _RGBColor(0xbb, 0x86, 0xfc)

    subtitle_box = s.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12.33), Inches(0.5))
    p = subtitle_box.text_frame.paragraphs[0]
    p.text = subtitle
    p.runs[0].font.size = Pt(16)
    p.runs[0].font.color.rgb = _RGBColor(0x94, 0xa3, 0xb8)

    return s


def add_roadmap_slide(prs):
    s = add_dark_slide(prs, ROADMAP["title"], ROADMAP["subtitle"])

    col_width = Inches(4.0)
    col_gap = Inches(0.2)
    left = Inches(0.5)
    top = Inches(1.9)
    height = Inches(5.3)

    for i, (ptitle, pbody) in enumerate(ROADMAP["phases"]):
        x = left + (col_width + col_gap) * i
        box = s.shapes.add_textbox(x, top, col_width, height)
        tf = box.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = ptitle
        p.runs[0].font.size = Pt(18)
        p.runs[0].font.bold = True
        p.runs[0].font.color.rgb = _RGBColor(0x7b, 0x68, 0xee)

        for line in pbody.split("\n"):
            pp = tf.add_paragraph()
            pp.text = line
            if pp.runs:
                pp.runs[0].font.size = Pt(13)
                pp.runs[0].font.color.rgb = _RGBColor(0xe2, 0xe8, 0xf0)

    footer = s.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(12.33), Inches(0.4))
    p = footer.text_frame.paragraphs[0]
    p.text = "Full roadmap: ROADMAP.md  ·  R1-01 through R4-04 covering desktop, mobile, enterprise, and research moonshots"
    p.runs[0].font.size = Pt(11)
    p.runs[0].font.italic = True
    p.runs[0].font.color.rgb = _RGBColor(0x94, 0xa3, 0xb8)

    # slide number
    num = s.shapes.add_textbox(Inches(12.8), Inches(7.1), Inches(0.3), Inches(0.3))
    p = num.text_frame.paragraphs[0]
    p.text = "14"
    p.runs[0].font.size = Pt(10)
    p.runs[0].font.color.rgb = _RGBColor(0x94, 0xa3, 0xb8)


def add_qa_slide(prs):
    s = add_dark_slide(prs, QA_SUMMARY_PLACEHOLDER["title"], QA_SUMMARY_PLACEHOLDER["subtitle"])

    stat_width = Inches(2.8)
    stat_height = Inches(2.0)
    stat_gap = Inches(0.3)
    left = Inches(0.7)
    top = Inches(2.0)

    for i, (num, label) in enumerate(QA_SUMMARY_PLACEHOLDER["stats"]):
        x = left + (stat_width + stat_gap) * i
        # big number
        box = s.shapes.add_textbox(x, top, stat_width, Inches(1.2))
        p = box.text_frame.paragraphs[0]
        p.text = num
        p.runs[0].font.size = Pt(60)
        p.runs[0].font.bold = True
        p.runs[0].font.color.rgb = _RGBColor(0xbb, 0x86, 0xfc)

        # label
        box = s.shapes.add_textbox(x, top + Inches(1.3), stat_width, Inches(0.5))
        p = box.text_frame.paragraphs[0]
        p.text = label
        p.runs[0].font.size = Pt(16)
        p.runs[0].font.color.rgb = _RGBColor(0x94, 0xa3, 0xb8)

    notes = s.shapes.add_textbox(Inches(0.7), Inches(4.8), Inches(11.9), Inches(2.0))
    tf = notes.text_frame
    tf.word_wrap = True
    for i, line in enumerate(QA_SUMMARY_PLACEHOLDER["notes"].split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        if p.runs:
            p.runs[0].font.size = Pt(14)
            p.runs[0].font.color.rgb = _RGBColor(0xe2, 0xe8, 0xf0)

    num = s.shapes.add_textbox(Inches(12.8), Inches(7.1), Inches(0.3), Inches(0.3))
    p = num.text_frame.paragraphs[0]
    p.text = "15"
    p.runs[0].font.size = Pt(10)
    p.runs[0].font.color.rgb = _RGBColor(0x94, 0xa3, 0xb8)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------


def main():
    root = Path(__file__).resolve().parent.parent
    src = root / "AccessBridge_Presentation.pptx"
    dst = root / "AccessBridge_Presentation_v2.pptx"

    prs = Presentation(src)
    print(f"Loaded {src.name} — {len(prs.slides)} slides")

    for i, pairs in REPLACEMENTS.items():
        print(f"Slide {i}: {len(pairs)} replacement(s)")
        replace_in_slide(prs.slides[i - 1], pairs)

    print("Appending Roadmap slide (14)")
    add_roadmap_slide(prs)
    print("Appending QA Summary slide (15)")
    add_qa_slide(prs)

    prs.save(dst)
    print(f"Wrote {dst.name} — {len(prs.slides)} slides total")


if __name__ == "__main__":
    main()
