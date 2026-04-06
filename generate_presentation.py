"""
AccessBridge — PowerPoint Presentation Generator v2
Compact, visually rich slides for Wipro TopGear Ideathon 2026.
"""

from pathlib import Path
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ── Theme ────────────────────────────────────────────────────────────────────
BG       = RGBColor(0x0F, 0x0F, 0x1E)
BG_CARD  = RGBColor(0x1A, 0x1A, 0x30)
ACCENT   = RGBColor(0x7B, 0x68, 0xEE)
ACCENT2  = RGBColor(0xBB, 0x86, 0xFC)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT    = RGBColor(0xD8, 0xD8, 0xE8)
MUTED    = RGBColor(0x90, 0x90, 0xA8)
GREEN    = RGBColor(0x10, 0xB9, 0x81)
ORANGE   = RGBColor(0xF5, 0x9E, 0x0B)
RED_SOFT = RGBColor(0xEF, 0x44, 0x44)
FONT     = "Calibri"

OUTPUT = Path(r"E:/code/AccessBridge/AccessBridge_Presentation.pptx")
SW = Emu(round(13.333 * 914400))  # slide width
SH = Inches(7.5)                   # slide height


# ── Helpers ──────────────────────────────────────────────────────────────────

def style(run, size, color, bold=False, italic=False):
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color

def add_bg(slide, prs):
    """Dark background + thin accent bar at bottom."""
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    r.fill.solid(); r.fill.fore_color.rgb = BG; r.line.fill.background()
    b = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, prs.slide_height - Inches(0.06), prs.slide_width, Inches(0.06))
    b.fill.solid(); b.fill.fore_color.rgb = ACCENT; b.line.fill.background()

def add_slide_num(slide, prs, n):
    bx = slide.shapes.add_textbox(prs.slide_width - Inches(0.8), prs.slide_height - Inches(0.38), Inches(0.6), Inches(0.25))
    tf = bx.text_frame; tf.clear()
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
    r = p.add_run(); r.text = str(n); style(r, 9, MUTED)

def add_title(slide, text, y=0.3, size=26):
    """Compact title with accent underline."""
    bx = slide.shapes.add_textbox(Inches(0.6), Inches(y), Inches(12.0), Inches(0.55))
    tf = bx.text_frame; tf.clear(); tf.word_wrap = True; tf.auto_size = MSO_AUTO_SIZE.NONE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = text; style(r, size, WHITE, bold=True)
    ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(y + 0.52), Inches(1.8), Inches(0.035))
    ln.fill.solid(); ln.fill.fore_color.rgb = ACCENT; ln.line.fill.background()

def add_subtitle_text(slide, text, y=0.88):
    bx = slide.shapes.add_textbox(Inches(0.6), Inches(y), Inches(12.0), Inches(0.35))
    tf = bx.text_frame; tf.clear(); tf.word_wrap = True
    p = tf.paragraphs[0]; r = p.add_run(); r.text = text; style(r, 12, MUTED, italic=True)

def add_card(slide, x, y, w, h, title_text, body_lines, title_color=ACCENT, body_size=11):
    """Rounded-corner card with title + bullet lines."""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.solid(); card.fill.fore_color.rgb = BG_CARD; card.line.fill.background()
    # Accent left border
    lb = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y + 0.08), Inches(0.04), Inches(h - 0.16))
    lb.fill.solid(); lb.fill.fore_color.rgb = title_color; lb.line.fill.background()

    bx = slide.shapes.add_textbox(Inches(x + 0.18), Inches(y + 0.1), Inches(w - 0.3), Inches(h - 0.18))
    tf = bx.text_frame; tf.clear(); tf.word_wrap = True; tf.auto_size = MSO_AUTO_SIZE.NONE
    # Title
    p = tf.paragraphs[0]; p.space_after = Pt(3)
    r = p.add_run(); r.text = title_text; style(r, 13, title_color, bold=True)
    # Body lines
    for line in body_lines:
        p = tf.add_paragraph(); p.space_before = Pt(1); p.space_after = Pt(1); p.line_spacing = 1.05
        r = p.add_run(); r.text = line; style(r, body_size, LIGHT)

def add_stat_box(slide, x, y, w, h, value, label, color=ACCENT):
    """Stat highlight box with large number."""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.solid(); card.fill.fore_color.rgb = BG_CARD; card.line.fill.background()
    # Value
    bx = slide.shapes.add_textbox(Inches(x + 0.1), Inches(y + 0.08), Inches(w - 0.2), Inches(h * 0.55))
    tf = bx.text_frame; tf.clear()
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = value; style(r, 22, color, bold=True)
    # Label
    bx2 = slide.shapes.add_textbox(Inches(x + 0.1), Inches(y + h * 0.5), Inches(w - 0.2), Inches(h * 0.45))
    tf2 = bx2.text_frame; tf2.clear()
    p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = label; style(r2, 9, MUTED)

def add_bullets(slide, x, y, w, h, lines, size=12, spacing=1.1, marker=""):
    """Compact bullet list."""
    bx = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = bx.text_frame; tf.clear(); tf.word_wrap = True; tf.auto_size = MSO_AUTO_SIZE.NONE
    for i, (text, bold_flag) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(3); p.space_before = Pt(1); p.line_spacing = spacing
        if marker:
            mr = p.add_run(); mr.text = marker + " "; style(mr, size, ACCENT, bold=True)
        r = p.add_run(); r.text = text
        style(r, size, WHITE if bold_flag else LIGHT, bold=bold_flag)

def add_two_col_bullets(slide, x, y, w, h, left_title, left_items, right_title, right_items):
    """Two-column bullet layout with headers."""
    half = w / 2 - 0.15
    # Left column
    add_card(slide, x, y, half, h, left_title, left_items)
    # Right column
    add_card(slide, x + half + 0.3, y, half, h, right_title, right_items)


# ── Slide Builders ───────────────────────────────────────────────────────────

def slide_01_title(slide, prs):
    """Title slide — centered, impactful."""
    # Main title
    bx = slide.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(11.33), Inches(1.4))
    tf = bx.text_frame; tf.clear(); tf.word_wrap = True; tf.auto_size = MSO_AUTO_SIZE.NONE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "AccessBridge"; style(r, 54, WHITE, bold=True)
    # Accent line
    lw = Inches(5)
    lx = (prs.slide_width - lw) // 2
    ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, lx, Inches(2.9), lw, Inches(0.04))
    ln.fill.solid(); ln.fill.fore_color.rgb = ACCENT; ln.line.fill.background()
    # Tagline
    bx2 = slide.shapes.add_textbox(Inches(1.5), Inches(3.15), Inches(10.33), Inches(0.6))
    tf2 = bx2.text_frame; tf2.clear(); tf2.word_wrap = True
    p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = "Ambient Accessibility Operating Layer"; style(r2, 22, ACCENT2, bold=True)
    # Sub-info
    bx3 = slide.shapes.add_textbox(Inches(2.0), Inches(3.85), Inches(9.33), Inches(0.8))
    tf3 = bx3.text_frame; tf3.clear(); tf3.word_wrap = True
    p3 = tf3.paragraphs[0]; p3.alignment = PP_ALIGN.CENTER
    r3 = p3.add_run(); r3.text = "Zero disclosure  |  AI-powered  |  Cross-application  |  Privacy-first"; style(r3, 14, MUTED)
    p4 = tf3.add_paragraph(); p4.alignment = PP_ALIGN.CENTER; p4.space_before = Pt(14)
    r4 = p4.add_run(); r4.text = "Wipro TopGear Ideathon 2026"; style(r4, 16, LIGHT)
    # Stat boxes across bottom
    stats = [("10+", "Features"), ("25+", "Voice Cmds"), ("116", "Unit Tests"), ("3-Tier", "AI Engine"), ("0", "Server Deps")]
    sx = 1.8
    for val, lbl in stats:
        add_stat_box(slide, sx, 5.4, 1.7, 0.9, val, lbl)
        sx += 1.95

def slide_02_problem(slide, prs):
    add_title(slide, "The Problem")
    add_subtitle_text(slide, "Why current accessibility tools fail 3.8 billion people")
    # Left: key stats
    add_stat_box(slide, 0.6, 1.45, 2.6, 1.1, "1.3B", "People with Disabilities", RED_SOFT)
    add_stat_box(slide, 0.6, 2.7, 2.6, 1.1, "2.5B", "Aging Population", ORANGE)
    add_stat_box(slide, 0.6, 3.95, 2.6, 1.1, "71%", "Abandon Inaccessible Sites", RED_SOFT)
    # Right: problem bullets
    problems = [
        ("Current tools require users to disclose their disability", False),
        ("Solutions are app-specific — no universal coverage", False),
        ("Static tools don't adapt to individual needs or context", False),
        ("Stigma prevents adoption: users avoid 'accessibility settings'", False),
        ("Enterprise apps (banking, insurance) worst offenders", False),
        ("Cost of inaccessibility: $6.9B/yr in lost revenue (US alone)", True),
    ]
    add_bullets(slide, 3.5, 1.45, 9.0, 3.8, problems, size=13, marker="\u25b8")

def slide_03_solution(slide, prs):
    add_title(slide, "Our Solution")
    add_subtitle_text(slide, "AccessBridge: invisible, adaptive, universal accessibility")
    # 3 key pillars as cards
    add_card(slide, 0.6, 1.5, 3.8, 2.5, "Zero Disclosure", [
        "No settings page, no self-identification",
        "Silently observes interaction patterns",
        "Adapts without user ever knowing",
        "Removes stigma barrier completely",
    ], ACCENT)
    add_card(slide, 4.6, 1.5, 3.8, 2.5, "Universal Coverage", [
        "Works on ANY website automatically",
        "Gmail, Outlook, banking, insurance, etc.",
        "Chrome extension = zero integration",
        "Enterprise-ready from day one",
    ], GREEN)
    add_card(slide, 8.6, 1.5, 3.8, 2.5, "AI-Adaptive", [
        "Real-time behavioral signal analysis",
        "3-tier AI: local \u2192 Gemini \u2192 Claude",
        "Learns user patterns over time",
        "All processing on-device (privacy-first)",
    ], ACCENT2)
    # Bottom: comparison strip
    bx = slide.shapes.add_textbox(Inches(0.6), 4.25 * 914400, Inches(12.0), Inches(0.3))
    tf = bx.text_frame; tf.clear(); tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "Traditional: User configures \u2192 Tool applies  |  AccessBridge: Tool observes \u2192 Auto-adapts \u2192 User benefits"
    style(r, 11, MUTED, italic=True)

def slide_04_architecture(slide, prs):
    add_title(slide, "Architecture")
    add_subtitle_text(slide, "Monorepo: 3 packages, event-driven pipeline, encrypted local storage")
    # Three package cards
    add_card(slide, 0.6, 1.5, 3.8, 2.0, "@accessbridge/core", [
        "TypeScript types & interfaces",
        "ProfileStore (encrypted IndexedDB)",
        "StruggleDetector (10 signal types)",
        "DecisionEngine (11 adaptive rules)",
        "Zero external dependencies",
    ], ACCENT)
    add_card(slide, 4.6, 1.5, 3.8, 2.0, "@accessbridge/extension", [
        "Chrome Manifest V3 extension",
        "React popup (5 tabs, live polling)",
        "Content scripts (all-site injection)",
        "Side panel dashboard (real-time)",
        "Background service worker (AI host)",
    ], GREEN)
    add_card(slide, 8.6, 1.5, 3.8, 2.0, "@accessbridge/ai-engine", [
        "3-tier: local \u2192 Gemini \u2192 Claude",
        "Response caching + deduplication",
        "Cost tracker with daily budgets",
        "Summarizer + Simplifier services",
        "Sub-100ms local tier response",
    ], ACCENT2)
    # Pipeline flow
    bx = slide.shapes.add_textbox(Inches(0.6), Inches(3.7), Inches(12.0), Inches(0.35))
    tf = bx.text_frame; tf.clear(); tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "Signal Pipeline:  User Interaction \u2192 Signal Collection \u2192 StruggleDetector \u2192 DecisionEngine \u2192 Auto-Adaptation"
    style(r, 12, ACCENT, bold=True)
    # Tech details
    details = [
        ("Build: Vite + TypeScript strict, output 132KB content + 28KB background + 19KB sidepanel", False),
        ("Storage: Encrypted IndexedDB profiles, chrome.storage.local for state, zero network for a11y data", False),
        ("Deployment: pnpm workspaces, VPS Docker (API:8100, Observatory:8200, Nginx:8300)", False),
    ]
    add_bullets(slide, 0.6, 4.2, 12.0, 1.5, details, size=11, marker="\u25aa")

def slide_05_struggle(slide, prs):
    add_title(slide, "Key Innovation: Struggle Detection Engine")
    add_subtitle_text(slide, "10 behavioral signals \u2192 real-time struggle score \u2192 automatic adaptation")
    # Signal grid (2 rows x 5 cols)
    signals = [
        ("Typing Speed", "WPM drops & rhythm"),
        ("Mouse Jitter", "Cursor instability"),
        ("Scroll Velocity", "Rapid/erratic scroll"),
        ("Error Rate", "Misclicks & typos"),
        ("Hesitation", "Long pauses on UI"),
        ("Backspace Ratio", "Frequent corrections"),
        ("Zoom Changes", "Font size struggles"),
        ("Voice Strain", "STT confidence drop"),
        ("Gaze Patterns", "Eye tracking drift"),
        ("Fatigue Index", "Session degradation"),
    ]
    sx, sy = 0.6, 1.5
    for i, (name, desc) in enumerate(signals):
        col = i % 5
        row = i // 5
        cx = sx + col * 2.45
        cy = sy + row * 1.15
        add_stat_box(slide, cx, cy, 2.25, 1.0, name, desc, ACCENT if row == 0 else ACCENT2)
    # Decision engine info
    add_card(slide, 0.6, 3.95, 5.8, 1.6, "Decision Engine (11 Rules)", [
        "Each rule maps signal patterns to specific adaptations",
        "Confidence scoring prevents false positives (threshold: 0.6)",
        "Profile-aware: adapts differently per disability type",
        "Baseline learning: adjusts to individual interaction norms",
        "Reversible: all adaptations can be undone instantly",
    ], ACCENT)
    add_card(slide, 6.6, 3.95, 5.8, 1.6, "Struggle Score (0-100)", [
        "0-30: Normal interaction \u2192 no intervention",
        "30-50: Mild struggle \u2192 subtle assists (larger targets)",
        "50-70: Moderate \u2192 active help (simplify, voice, focus)",
        "70-90: High \u2192 aggressive adaptation (reflow, summarize)",
        "90-100: Critical \u2192 maximum assist (essential content only)",
    ], ORANGE)

def slide_06_cognitive(slide, prs):
    add_title(slide, "Cognitive Accessibility")
    add_subtitle_text(slide, "Focus, comprehension, and information overload support")
    # Feature cards
    add_card(slide, 0.6, 1.5, 4.0, 1.85, "Focus Mode", [
        "Spotlight follows cursor, periphery dims to 30%",
        "Removes sidebar ads, popups, notifications",
        "Auto-triggers on rapid tab switching or scroll jitter",
    ], ACCENT)
    add_card(slide, 4.8, 1.5, 4.0, 1.85, "Distraction Shield", [
        "Hides animations, auto-play videos, marquees",
        "Reduces visual noise by 60-80%",
        "Configurable strictness levels (mild/strong)",
    ], GREEN)
    add_card(slide, 9.0, 1.5, 3.4, 1.85, "Reading Guide", [
        "Ruler line tracks reading position",
        "Adjustable width, color, opacity",
        "Keyboard-driven (arrow keys)",
    ], ACCENT2)
    add_card(slide, 0.6, 3.55, 6.0, 2.0, "AI-Powered Summarization", [
        "Page summarization: extracts key bullets + reading time + complexity score",
        "Email summarization: Gmail/Outlook toolbar buttons auto-injected",
        "Slide-in summary panel with Read Aloud (Web Speech TTS), Copy, Dismiss",
        "Auto-summarize mode: triggers 2s after page load on complex pages",
        "MutationObserver handles SPA navigation (Gmail, Outlook dynamic loading)",
    ], ACCENT)
    add_card(slide, 6.8, 3.55, 5.6, 2.0, "Text Simplification", [
        "Flesch-Kincaid readability targeting (grade level reduction)",
        "Word-level: replaces complex terms with simpler synonyms",
        "Sentence-level: shortens and restructures for clarity",
        "Two modes: mild (gentle rewording) / strong (plain language)",
        "Preserves technical terms in context (banking, medical, legal)",
    ], ACCENT2)

def slide_07_motor(slide, prs):
    add_title(slide, "Motor Accessibility")
    add_subtitle_text(slide, "Voice, gaze, dwell, keyboard, and predictive input for zero-mouse operation")
    # Cards
    add_card(slide, 0.6, 1.5, 4.0, 2.1, "Voice Commands (25+)", [
        "Navigation: scroll up/down, go back/forward, top/bottom",
        "Actions: click [target], type [text], find [text]",
        "Tabs: next/prev tab, close tab, new tab",
        "A11y: zoom in/out, focus mode, read aloud",
        "Hindi: 'neeche scroll karo', 'vapas jao', etc.",
        "Language auto-select from user profile",
    ], ACCENT)
    add_card(slide, 4.8, 1.5, 3.8, 2.1, "Eye Tracking", [
        "FaceDetector API (zero deps, native Chrome)",
        "60% head pose + 40% eye offset blend",
        "Smooth gaze cursor with EMA filter",
        "5-point calibration system",
        "Skin-color centroid fallback",
        "Webcam preview + toggle",
    ], GREEN)
    add_card(slide, 8.8, 1.5, 3.6, 2.1, "Dwell Click", [
        "Radial SVG progress indicator",
        "Configurable delay (200-2000ms)",
        "15px movement threshold",
        "Visual pulse on click",
        "Target highlight on hover",
        "Auto-disable on mouse movement",
    ], ORANGE)
    add_card(slide, 0.6, 3.8, 5.8, 1.8, "Keyboard-Only Mode", [
        "Skip links: jump to main, nav, footer with one key",
        "Enhanced focus ring: visible 3px outline on all elements",
        "Tab order optimizer: auto-adds tabindex to clickable elements",
        "Shortcuts overlay (press '?'): shows all available hotkeys",
        "Arrow key navigation within groups, Escape to deselect",
        "MutationObserver: handles dynamic content (SPAs, infinite scroll)",
    ], ACCENT2)
    add_card(slide, 6.6, 3.8, 5.8, 1.8, "Predictive Input", [
        "Frequency-based word prediction (~500 word dictionary)",
        "Session learning: adapts to current typing patterns",
        "Floating suggestion panel: Alt+1-5 or Tab to accept",
        "Phrase auto-complete: ~50 common phrases",
        "Form field intelligence: detects email/phone/address/name",
        "80ms debounced, contenteditable support",
    ], ACCENT)

def slide_08_sensory_fatigue(slide, prs):
    add_title(slide, "Sensory Adaptation & Fatigue-Adaptive UI")
    add_subtitle_text(slide, "Visual accessibility + progressive UI simplification based on session degradation")
    # Left: Sensory
    add_card(slide, 0.6, 1.5, 5.8, 2.3, "Sensory Adaptation", [
        "High contrast: 4 themes (dark, light, yellow-on-black, custom)",
        "Font scaling: 50%-300% dynamic range, all DOM elements",
        "Text-to-speech: Web Speech API with rate/pitch/voice controls",
        "Reduced motion: disables CSS animations, transitions, auto-play",
        "Color filters: protanopia, deuteranopia, tritanopia, achromatopsia",
        "Profile-persistent: settings saved to encrypted IndexedDB",
    ], ACCENT)
    # Right: Fatigue levels
    add_card(slide, 6.6, 1.5, 5.8, 2.3, "Fatigue-Adaptive UI (4 Levels)", [
        "L1 Mild: larger click targets (+20%), simpler navigation labels",
        "L2 Moderate: reduce visible content, increase spacing by 1.5x",
        "L3 High: auto-summarize long pages, simplify layout to single-col",
        "L4 Critical: essential content only, break reminders every 15 min",
        "Signals: session duration, time-of-day, declining click accuracy",
        "Smooth transitions between levels, manual override available",
    ], ORANGE)
    # Bottom: how it works
    bx = slide.shapes.add_textbox(Inches(0.6), Inches(4.0), Inches(12.0), Inches(0.6))
    tf = bx.text_frame; tf.clear(); tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "Fatigue Pipeline:  Session Timer + Interaction Decay + Time-of-Day \u2192 Fatigue Index (0-100) \u2192 Level Selection \u2192 Progressive Simplification"
    style(r, 11, ACCENT, bold=True)

def slide_09_domains(slide, prs):
    add_title(slide, "Domain Intelligence")
    add_subtitle_text(slide, "Context-aware connectors for banking, insurance, and extensible to any vertical")
    add_card(slide, 0.6, 1.5, 5.8, 3.2, "Banking Connector", [
        "Jargon Decoder: 25 terms (EMI, CIBIL, NEFT, RTGS, etc.)",
        "Form Assistance: auto-validates account numbers, IFSC, PAN",
        "Indian Numbering: displays amounts as Lakh/Crore/Arab",
        "Security Alerts: warns on sensitive actions (transfers, password)",
        "Transaction Simplifier: plain-English descriptions of entries",
        "Auto-detects banking domains (SBI, HDFC, ICICI, etc.)",
    ], GREEN)
    add_card(slide, 6.6, 1.5, 5.8, 3.2, "Insurance Connector", [
        "Policy Simplifier: rewrites jargon-heavy clauses in plain language",
        "Jargon Decoder: 35 terms (premium, deductible, copay, etc.)",
        "Comparison Helper: side-by-side policy feature comparison",
        "Claim Form Assistant: step-by-step guidance with validation",
        "Premium Calculator: explains premium breakdown components",
        "Auto-detects insurance domains and activates connector",
    ], ACCENT2)
    # Architecture note
    add_card(slide, 0.6, 4.9, 12.0, 0.7, "Extensible Architecture", [
        "DomainConnectorRegistry auto-detects domain from URL \u2192 activates matching connector \u2192 injects domain-specific UI + helpers. Add new connectors in <50 lines.",
    ], ACCENT, body_size=11)

def slide_10_ai(slide, prs):
    add_title(slide, "AI Engine: 3-Tier Cost Optimization")
    add_subtitle_text(slide, "Smart routing: free local processing first, cloud only when needed")
    # Tier cards
    add_card(slide, 0.6, 1.5, 3.8, 2.8, "Tier 1: Local (Free)", [
        "Extractive summarization (TF-IDF)",
        "Rule-based word simplification",
        "Regex classification + heuristics",
        "Response time: <100ms",
        "Cost: $0.00",
        "Handles ~80% of all requests",
    ], GREEN)
    add_card(slide, 4.6, 1.5, 3.8, 2.8, "Tier 2: Gemini Flash", [
        "Abstractive summarization",
        "Complex text simplification",
        "Context-aware classification",
        "Response time: 500-2000ms",
        "Cost: ~$0.001/request",
        "Handles ~15% of requests",
    ], ACCENT)
    add_card(slide, 8.6, 1.5, 3.8, 2.8, "Tier 3: Claude", [
        "Advanced reasoning tasks",
        "Multi-step simplification",
        "Nuanced accessibility analysis",
        "Response time: 1-5s",
        "Cost: ~$0.01/request",
        "Handles ~5% of requests",
    ], ACCENT2)
    # Infrastructure details
    add_card(slide, 0.6, 4.5, 6.0, 1.1, "Caching & Deduplication", [
        "SHA-256 key hashing, configurable TTL, hit/miss rate tracking",
        "Request dedup: identical in-flight requests share one API call",
    ], ACCENT)
    add_card(slide, 6.8, 4.5, 5.6, 1.1, "Cost Tracking & Budgets", [
        "Per-tier cost estimation, daily budget enforcement ($1/day default)",
        "Graceful degradation: falls back to lower tier when budget exhausted",
    ], ORANGE)

def slide_11_testing(slide, prs):
    add_title(slide, "Testing, Quality & Tech Stack")
    add_subtitle_text(slide, "116 tests, zero errors, strict TypeScript, automated deployment pipeline")
    # Stats row
    stats = [("116", "Tests Passing"), ("3", "Test Suites"), ("0", "TS Errors"), ("0", "Build Warnings"), ("100%", "Core Coverage")]
    sx = 0.6
    for val, lbl in stats:
        clr = GREEN if val in ("116", "100%") else ACCENT
        add_stat_box(slide, sx, 1.5, 2.2, 1.0, val, lbl, clr)
        sx += 2.4
    # Test details
    add_card(slide, 0.6, 2.7, 6.0, 1.6, "Test Coverage", [
        "StruggleDetector: 16 tests (signal processing, baseline, scoring)",
        "DecisionEngine: 21 tests (rules, confidence, revert, profiles)",
        "ProfileStore: 25 tests (CRUD, encryption, migration, edge cases)",
        "AI Engine: 54 tests (cache, normalizer, cost-tracker, local-provider)",
    ], GREEN)
    add_card(slide, 6.8, 2.7, 5.6, 1.6, "Tech Stack", [
        "Chrome Extension (Manifest V3) + React + TypeScript + Vite",
        "pnpm monorepo: @accessbridge/core + extension + ai-engine",
        "IndexedDB (encrypted), chrome.storage.local, Web Speech API",
        "FaceDetector API, vitest, Docker (VPS), Nginx reverse proxy",
    ], ACCENT)
    # CI/CD pipeline
    bx = slide.shapes.add_textbox(Inches(0.6), Inches(4.5), Inches(12.0), Inches(0.35))
    tf = bx.text_frame; tf.clear(); tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "CI/CD:  pnpm build \u2192 tsc --noEmit \u2192 vitest run \u2192 git push \u2192 SSH deploy \u2192 docker-compose up"
    style(r, 11, ACCENT, bold=True)

def slide_12_impact(slide, prs):
    add_title(slide, "Impact & Market Opportunity")
    add_subtitle_text(slide, "Serving 3.8 billion people who are currently underserved or excluded")
    # Market stats
    add_stat_box(slide, 0.6, 1.5, 2.8, 1.2, "3.8B", "Total Addressable Market", ACCENT)
    add_stat_box(slide, 3.6, 1.5, 2.8, 1.2, "$6.9B", "Lost Revenue/yr (US)", RED_SOFT)
    add_stat_box(slide, 6.6, 1.5, 2.8, 1.2, "71%", "Abandon Inaccessible Sites", ORANGE)
    add_stat_box(slide, 9.6, 1.5, 2.8, 1.2, "98%", "Top Sites Fail WCAG", RED_SOFT)
    # Value propositions
    add_card(slide, 0.6, 2.95, 6.0, 2.7, "Why AccessBridge Wins", [
        "Zero disclosure: removes the #1 adoption barrier (stigma)",
        "Universal: no per-app integration, works everywhere instantly",
        "Privacy-first: all behavioral data stays on-device",
        "AI-adaptive: learns and improves, not static rules",
        "Cost-optimized: 80% of AI requests are free (local tier)",
        "Enterprise-ready: banking + insurance connectors built-in",
    ], ACCENT)
    add_card(slide, 6.8, 2.95, 5.6, 2.7, "Business Model & Growth", [
        "B2C: Free Chrome extension, premium AI features",
        "B2B: Enterprise licensing for WCAG/ADA compliance",
        "TAM: 1.3B disabled + 2.5B aging + enterprise orgs",
        "Regulatory tailwinds: ADA, EAA, WCAG 2.2 mandates",
        "Network effect: user profiles improve AI accuracy",
        "Expansion: Firefox, Edge, mobile, native apps",
    ], GREEN)

def slide_13_demo(slide, prs):
    add_title(slide, "Live Demo & Contact")
    # Center content
    bx = slide.shapes.add_textbox(Inches(1.5), Inches(1.5), Inches(10.33), Inches(1.0))
    tf = bx.text_frame; tf.clear(); tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "Chrome Extension Demo Available"; style(r, 22, WHITE, bold=True)
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER; p2.space_before = Pt(10)
    r2 = p2.add_run(); r2.text = "Load packages/extension/dist/ in chrome://extensions"; style(r2, 14, MUTED)

    # Demo features grid
    demo_features = [
        ("Struggle Detection", "Live score 0-100"),
        ("Voice Commands", "25+ English + Hindi"),
        ("Focus Mode", "Spotlight + shield"),
        ("Eye Tracking", "Webcam gaze cursor"),
        ("Dwell Click", "Radial auto-click"),
        ("AI Summarize", "Email & page bullets"),
        ("Keyboard Mode", "Full keyboard nav"),
        ("Predictive Input", "Smart suggestions"),
        ("Domain Connectors", "Banking + Insurance"),
        ("Fatigue UI", "4-level simplification"),
    ]
    sx, sy = 0.8, 2.8
    for i, (feat, desc) in enumerate(demo_features):
        col = i % 5; row = i // 5
        cx = sx + col * 2.45; cy = sy + row * 1.15
        add_stat_box(slide, cx, cy, 2.25, 1.0, feat, desc, ACCENT if row == 0 else GREEN)

    # Links
    bx2 = slide.shapes.add_textbox(Inches(1.5), Inches(5.25), Inches(10.33), Inches(1.0))
    tf2 = bx2.text_frame; tf2.clear(); tf2.word_wrap = True
    p3 = tf2.paragraphs[0]; p3.alignment = PP_ALIGN.CENTER
    r3 = p3.add_run(); r3.text = "GitHub: github.com/manishjnv/AccessBridge"; style(r3, 14, ACCENT, bold=True)
    p4 = tf2.add_paragraph(); p4.alignment = PP_ALIGN.CENTER; p4.space_before = Pt(12)
    r4 = p4.add_run(); r4.text = "Team AccessBridge  |  Wipro TopGear Ideathon 2026"; style(r4, 13, LIGHT)
    p5 = tf2.add_paragraph(); p5.alignment = PP_ALIGN.CENTER; p5.space_before = Pt(16)
    r5 = p5.add_run(); r5.text = "Thank You!"; style(r5, 28, ACCENT2, bold=True)


# ── Build ────────────────────────────────────────────────────────────────────

BUILDERS = [
    slide_01_title,
    slide_02_problem,
    slide_03_solution,
    slide_04_architecture,
    slide_05_struggle,
    slide_06_cognitive,
    slide_07_motor,
    slide_08_sensory_fatigue,
    slide_09_domains,
    slide_10_ai,
    slide_11_testing,
    slide_12_impact,
    slide_13_demo,
]

def build():
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH
    blank = prs.slide_layouts[6]
    for i, builder in enumerate(BUILDERS):
        slide = prs.slides.add_slide(blank)
        add_bg(slide, prs)
        builder(slide, prs)
        add_slide_num(slide, prs, i + 1)
    prs.save(str(OUTPUT))
    print(f"Saved: {OUTPUT}  ({len(BUILDERS)} slides)")

if __name__ == "__main__":
    build()
